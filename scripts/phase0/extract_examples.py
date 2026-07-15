#!/usr/bin/env python3
"""只读提取 Example 样本，并生成本地 Manifest 与 Markdown 文本。"""

from __future__ import annotations

import argparse
import hashlib
import json
import mimetypes
import subprocess
import tempfile
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation


SUPPORTED_SUFFIXES = {".pdf", ".doc", ".docx", ".pptx", ".mp3"}


@dataclass(frozen=True)
class SourceRecord:
    """描述一个只读样本及其本地提取状态。"""

    source_id: str
    filename: str
    suffix: str
    media_type: str
    size_bytes: int
    sha256: str
    source_role: str
    confidentiality: str
    current_validity: str
    extraction_status: str
    extracted_markdown: str | None
    transcript_markdown: str | None
    notes: list[str]


def sha256_file(path: Path) -> str:
    """流式计算文件 SHA-256，避免把大文件一次性读入内存。"""

    digest = hashlib.sha256()
    with path.open("rb") as source:
        for chunk in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def infer_source_role(path: Path) -> str:
    """根据文件类型和名称给出保守的来源角色建议。"""

    name = path.name.lower()
    if path.suffix.lower() == ".mp3":
        return "meeting_audio"
    if "纪要" in name:
        return "meeting_minutes"
    if path.suffix.lower() == ".pptx":
        return "planning_deck"
    if "发票" in name:
        return "administrative_evidence"
    if path.suffix.lower() in {".doc", ".docx"}:
        return "working_document"
    return "reference_document"


def infer_confidentiality(source_role: str) -> str:
    """按来源角色应用最小披露等级，行政凭证默认视为 P3。"""

    if source_role == "administrative_evidence":
        return "P3"
    return "P2"


def extract_pdf(path: Path) -> str:
    """按页提取 PDF 文本与可识别表格。"""

    sections: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page_number, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()
            tables = page.extract_tables() or []
            body: list[str] = [f"## 第 {page_number} 页"]
            body.append(text or "_未提取到可搜索文本，可能需要本地 OCR。_")
            for table_number, table in enumerate(tables, start=1):
                body.append(f"\n### 表格 {table_number}\n")
                for row in table:
                    body.append(" | ".join((cell or "").replace("\n", " ") for cell in row))
            sections.append("\n".join(body))
    return "\n\n".join(sections)


def extract_docx(path: Path) -> str:
    """提取 DOCX 正文与表格，保持原始顺序的可读近似。"""

    document = Document(path)
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table_number, table in enumerate(document.tables, start=1):
        parts.append(f"\n## 表格 {table_number}")
        for row in table.rows:
            parts.append(" | ".join(cell.text.replace("\n", " ").strip() for cell in row.cells))
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    """按幻灯片提取可见文本和演讲者备注。"""

    presentation = Presentation(path)
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    texts.append(value)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    texts.append(" | ".join(cell.text.replace("\n", " ").strip() for cell in row.cells))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except (AttributeError, ValueError):
            notes = ""
        if notes:
            texts.append(f"\n### 演讲者备注\n{notes}")
        slides.append(f"## 第 {index} 页\n\n" + ("\n\n".join(texts) or "_无可见文本。_"))
    return "\n\n".join(slides)


def extract_doc(path: Path) -> str:
    """使用 macOS textutil 在临时目录中只读转换旧版 DOC。"""

    with tempfile.TemporaryDirectory(prefix="brand-os-doc-") as temp_dir:
        output = Path(temp_dir) / "converted.txt"
        result = subprocess.run(
            ["/usr/bin/textutil", "-convert", "txt", "-output", str(output), str(path)],
            check=False,
            capture_output=True,
            text=True,
        )
        if result.returncode != 0 or not output.exists():
            detail = (result.stderr or result.stdout).strip()
            raise RuntimeError(f"textutil 转换失败：{detail or result.returncode}")
        return output.read_text(encoding="utf-8", errors="replace").strip()


def probe_audio(path: Path) -> tuple[str, dict[str, object]]:
    """读取音频元数据；转写由独立的本地流程完成。"""

    result = subprocess.run(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_format",
            "-show_streams",
            "-of",
            "json",
            str(path),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    metadata = json.loads(result.stdout)
    duration = metadata.get("format", {}).get("duration", "unknown")
    markdown = f"## 音频元数据\n\n- 时长（秒）：{duration}\n- 状态：待本地转写"
    return markdown, metadata


def extract_content(path: Path) -> tuple[str, str, list[str]]:
    """提取一个受支持文件，返回 Markdown、状态和备注。"""

    suffix = path.suffix.lower()
    notes: list[str] = []
    if suffix == ".pdf":
        content = extract_pdf(path)
        if "可能需要本地 OCR" in content:
            notes.append("PDF 未获得完整可搜索文本，需要 OCR 复核")
        return content, "extracted", notes
    if suffix == ".docx":
        return extract_docx(path), "extracted", notes
    if suffix == ".doc":
        return extract_doc(path), "extracted", notes
    if suffix == ".pptx":
        return extract_pptx(path), "extracted", notes
    if suffix == ".mp3":
        content, _ = probe_audio(path)
        notes.append("仅提取音频元数据，正文等待本地转写")
        return content, "metadata_only", notes
    raise ValueError(f"不支持的文件类型：{suffix}")


def ensure_within_root(path: Path, root: Path) -> None:
    """阻止符号链接或路径逃逸到授权样本根目录之外。"""

    resolved = path.resolve(strict=True)
    root_resolved = root.resolve(strict=True)
    if not resolved.is_relative_to(root_resolved):
        raise ValueError(f"文件越出授权目录：{path}")


def build_manifest(source_root: Path, output_root: Path) -> dict[str, object]:
    """提取目录内全部受支持文件，并返回本地 Manifest。"""

    source_root = source_root.resolve(strict=True)
    output_root.mkdir(parents=True, exist_ok=True)
    extracted_root = output_root / "extracted"
    extracted_root.mkdir(parents=True, exist_ok=True)

    records: list[SourceRecord] = []
    seen_hashes: dict[str, str] = {}
    for path in sorted(source_root.iterdir(), key=lambda item: item.name):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_SUFFIXES:
            continue
        ensure_within_root(path, source_root)
        digest = sha256_file(path)
        source_id = f"EX-{digest[:12].upper()}"
        markdown_path = extracted_root / f"{source_id}.md"
        transcript_path = output_root / "audio" / f"{source_id}.md"
        notes: list[str] = []
        try:
            content, status, extraction_notes = extract_content(path)
            notes.extend(extraction_notes)
            markdown_path.write_text(
                f"# {source_id}\n\n> 本地只读提取，原件 SHA-256：`{digest}`\n\n{content}\n",
                encoding="utf-8",
            )
            extracted_reference: str | None = str(markdown_path.relative_to(output_root))
        except Exception as error:  # 保留单文件失败，避免吞掉批次其他证据
            status = "failed"
            extracted_reference = None
            notes.append(str(error))
        if digest in seen_hashes:
            notes.append(f"内容与 {seen_hashes[digest]} 重复")
        else:
            seen_hashes[digest] = source_id
        media_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
        source_role = infer_source_role(path)
        transcript_reference = (
            str(transcript_path.relative_to(output_root)) if transcript_path.exists() else None
        )
        if path.suffix.lower() == ".mp3" and transcript_reference:
            notes.append("已有本地模型候选转写，未经人工校对")
        records.append(
            SourceRecord(
                source_id=source_id,
                filename=path.name,
                suffix=path.suffix.lower(),
                media_type=media_type,
                size_bytes=path.stat().st_size,
                sha256=digest,
                source_role=source_role,
                confidentiality=infer_confidentiality(source_role),
                current_validity="unknown",
                extraction_status=status,
                extracted_markdown=extracted_reference,
                transcript_markdown=transcript_reference,
                notes=notes,
            )
        )

    manifest = {
        "schema_version": "sample-manifest.v1",
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "source_root": str(source_root),
        "scope": "example_only",
        "authority": "local_read_only_user_provided",
        "record_count": len(records),
        "known_source_gap": True,
        "records": [asdict(record) for record in records],
    }
    (output_root / "manifest.local.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return manifest


def parse_args() -> argparse.Namespace:
    """解析命令行参数。"""

    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("source_root", type=Path, help="获授权的只读样本目录")
    parser.add_argument("--output", type=Path, required=True, help="本地解析输出目录")
    return parser.parse_args()


def main() -> int:
    """执行本地提取并输出简要结果。"""

    args = parse_args()
    manifest = build_manifest(args.source_root, args.output)
    print(json.dumps({"record_count": manifest["record_count"], "output": str(args.output)}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
