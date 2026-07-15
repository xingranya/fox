"""Phase 0 本地样本提取器测试。"""

from __future__ import annotations

import importlib.util
import sys
import tempfile
import unittest
from pathlib import Path

from docx import Document
from pptx import Presentation


MODULE_PATH = Path(__file__).parents[2] / "scripts" / "phase0" / "extract_examples.py"
SPEC = importlib.util.spec_from_file_location("extract_examples", MODULE_PATH)
assert SPEC and SPEC.loader
extract_examples = importlib.util.module_from_spec(SPEC)
sys.modules[SPEC.name] = extract_examples
SPEC.loader.exec_module(extract_examples)


class ExtractExamplesTest(unittest.TestCase):
    """验证稳定 ID、角色识别、Office 提取和目录边界。"""

    def test_sha256_is_stable(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "sample.bin"
            path.write_bytes(b"brand-os")
            self.assertEqual(
                extract_examples.sha256_file(path),
                "779d8fc9df09cb10c4e18b6c36f29794d228c45ad55eff28aab990e4e46bf320",
            )

    def test_infer_source_role_is_conservative(self) -> None:
        self.assertEqual(extract_examples.infer_source_role(Path("会议纪要.docx")), "meeting_minutes")
        self.assertEqual(extract_examples.infer_source_role(Path("规划.pptx")), "planning_deck")
        self.assertEqual(extract_examples.infer_source_role(Path("录音.mp3")), "meeting_audio")
        self.assertEqual(extract_examples.infer_confidentiality("administrative_evidence"), "P3")
        self.assertEqual(extract_examples.infer_confidentiality("meeting_minutes"), "P2")

    def test_extract_docx_and_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            docx_path = root / "纪要.docx"
            document = Document()
            document.add_paragraph("这是会议内容")
            document.save(docx_path)

            pptx_path = root / "规划.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "半年度规划"
            presentation.save(pptx_path)

            self.assertIn("这是会议内容", extract_examples.extract_docx(docx_path))
            self.assertIn("半年度规划", extract_examples.extract_pptx(pptx_path))

    def test_build_manifest_keeps_source_read_only(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir) / "source"
            output = Path(temp_dir) / "output"
            root.mkdir()
            document = Document()
            document.add_paragraph("只读样本")
            source = root / "样本.docx"
            document.save(source)
            before = source.read_bytes()

            manifest = extract_examples.build_manifest(root, output)

            self.assertEqual(manifest["record_count"], 1)
            self.assertEqual(source.read_bytes(), before)
            record = manifest["records"][0]
            self.assertEqual(record["extraction_status"], "extracted")
            self.assertTrue((output / record["extracted_markdown"]).exists())


if __name__ == "__main__":
    unittest.main()
